"""Prepend a team slide to the NotebookLM deck.

The deck is one full-bleed image per slide on a 17.78x10in canvas, background #171717
with a teal accent and a faint blueprint grid. We rebuild that look with native shapes
so the new slide is editable and matches, then move it to position 0.
"""

import copy
import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

BG = RGBColor(0x17, 0x17, 0x17)
GRID = RGBColor(0x24, 0x2A, 0x2C)
TEAL = RGBColor(0x7F, 0xB3, 0xC8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREY = RGBColor(0x9A, 0x9A, 0x9A)
AMBER = RGBColor(0xD2, 0x99, 0x4A)
MONO = "Consolas"
SANS = "Segoe UI"

MEMBERS = [
    ("ARAVINDA KANNAN KS", "24BCE1290", True),
    ("TARUN KRISHNA MANIVANNAN", "24BCE5460", False),
    ("KRISHNA PRASAD M", "24BCE5519", False),
    ("ELUKE. SRINITHA", "24BCE1639", False),
]


def textbox(slide, x, y, w, h, text, *, size, color, font=SANS, bold=False,
            align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.name = font
        r.font.bold = bold
    return tb


def rect(slide, x, y, w, h, *, fill=None, line=None, width=1.0, shape=MSO_SHAPE.RECTANGLE):
    sh = slide.shapes.add_shape(shape, x, y, w, h)
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(width)
    sh.shadow.inherit = False
    return sh


def build(path_in: str, path_out: str) -> None:
    prs = Presentation(path_in)
    W, H = prs.slide_width, prs.slide_height
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # background
    rect(slide, 0, 0, W, H, fill=BG)

    # faint blueprint grid
    step = Inches(0.78)
    x = 0
    while x < W:
        rect(slide, x, 0, Emu(9525), H, fill=GRID)
        x += step
    y = 0
    while y < H:
        rect(slide, 0, y, W, Emu(9525), fill=GRID)
        y += step

    # outer frame
    rect(slide, Inches(0.42), Inches(0.34), W - Inches(0.84), H - Inches(0.68),
         line=WHITE, width=1.1)

    # top-right track chip
    chip_w, chip_h = Inches(6.1), Inches(1.12)
    cx, cy = W - Inches(0.42) - chip_w, Inches(0.34)
    rect(slide, cx, cy, chip_w, chip_h, line=WHITE, width=1.1)
    textbox(slide, cx + Inches(0.3), cy + Inches(0.13), chip_w - Inches(0.6), Inches(0.4),
            "Theme A — FINSWARM", size=19, color=TEAL, font=MONO, align=PP_ALIGN.RIGHT)
    textbox(slide, cx + Inches(0.3), cy + Inches(0.58), chip_w - Inches(0.6), Inches(0.34),
            "Agentic Swarm · The AI Boardroom", size=13, color=GREY, font=MONO,
            align=PP_ALIGN.RIGHT)

    # eyebrow
    textbox(slide, Inches(1.15), Inches(2.15), Inches(8), Inches(0.4),
            "T  E  A  M", size=15, color=TEAL, font=MONO)

    # team name
    textbox(slide, Inches(1.1), Inches(2.6), W - Inches(2.2), Inches(1.5),
            "404 : FUNDS NOT FOUND", size=68, color=WHITE, font=MONO, bold=True)

    # teal rule
    rect(slide, Inches(1.15), Inches(4.28), Inches(4.6), Emu(28575), fill=TEAL)

    # members
    top = Inches(4.95)
    row_h = Inches(0.92)
    for i, (name, reg, lead) in enumerate(MEMBERS):
        y = top + row_h * i
        rect(slide, Inches(1.15), y + Inches(0.16), Emu(28575), Inches(0.46),
             fill=AMBER if lead else GRID)
        textbox(slide, Inches(1.45), y, Inches(8.6), Inches(0.5), name,
                size=22, color=WHITE, font=SANS, bold=lead)
        textbox(slide, Inches(1.45), y + Inches(0.46), Inches(8.6), Inches(0.34), reg,
                size=15, color=GREY, font=MONO)
        if lead:
            tag_w = Inches(1.72)
            tag = rect(slide, Inches(1.45) + Inches(3.62), y + Inches(0.02),
                       tag_w, Inches(0.38), line=AMBER, width=1.0)
            textbox(slide, Inches(1.45) + Inches(3.62), y + Inches(0.02), tag_w, Inches(0.38),
                    "TEAM LEAD", size=12, color=AMBER, font=MONO, align=PP_ALIGN.CENTER)

    # right-hand strapline
    bx = Inches(11.2)
    rect(slide, bx, Inches(4.95), W - bx - Inches(1.0), Inches(3.4), line=GREY, width=0.9)
    textbox(slide, bx + Inches(0.42), Inches(5.42), W - bx - Inches(1.85), Inches(1.5),
            "A board that convenes\non any business problem.", size=25, color=WHITE, font=SANS)
    textbox(slide, bx + Inches(0.42), Inches(7.05), W - bx - Inches(1.85), Inches(0.9),
            "Six mandated agents. Hard constraints\nchecked, not weighted.",
            size=14, color=GREY, font=MONO)

    # move the new slide to the front
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    sldIdLst.remove(ids[-1])
    sldIdLst.insert(0, ids[-1])

    prs.save(path_out)
    print("saved", path_out)


if __name__ == "__main__":
    build(sys.argv[1], sys.argv[2])
